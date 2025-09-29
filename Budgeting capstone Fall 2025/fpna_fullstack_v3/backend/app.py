"""
FP&A Fullstack Backend (FastAPI) — v3
New in v3:
- Headcount/Roster module → translates hires/raises into Labor/Taxes/Benefits plan rows
- Goal Seek solver → hit a target EBITDA by adjusting chosen levers (revenue %, cogs %, labor %)
- All v2 features: Auth/Roles, Budgets, Scenarios, GDP presets, AI forecasts, BvA, Board pack, QB import, Excel endpoints

Run: uvicorn app:api --reload --port 8000
"""
from __future__ import annotations
import os, io, csv, math
from typing import Dict, Optional, List, Literal, Tuple
import numpy as np
import pandas as pd

from fastapi import FastAPI, Body, HTTPException, Depends, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials

from statsmodels.tsa.arima.model import ARIMA
from sklearn.linear_model import LinearRegression

from pptx import Presentation
from pptx.util import Inches

security = HTTPBearer(auto_error=False)
TOKENS = {
    "admin-token": {"user":"alice","role":"Admin"},
    "analyst-token": {"user":"bob","role":"Analyst"},
    "viewer-token": {"user":"cathy","role":"Viewer"},
}
def get_current_user(creds: Optional[HTTPAuthorizationCredentials] = Depends(security)):
    if creds is None or not creds.credentials:
        raise HTTPException(401, "Missing Bearer token")
    info = TOKENS.get(creds.credentials)
    if not info:
        raise HTTPException(403, "Invalid token")
    return info
def require_role(*roles: str):
    def checker(user=Depends(get_current_user)):
        if user["role"] not in roles:
            raise HTTPException(403, f"Requires role: {', '.join(roles)}")
        return user
    return checker

FUNC_CCY = os.getenv("FUNC_CURRENCY", "USD")
YEARS = [2024, 2025, 2026, 2027]
MONTHS = [f"{m:02d}" for m in range(1,13)]
ACCOUNTS = [
    "Revenue:Food","Revenue:Beverage",
    "COGS:Food","COGS:Paper",
    "Opex:Labor","Opex:PayrollTaxes","Opex:Benefits",
    "Opex:Rent","Opex:Utilities","Opex:Royalty","Opex:AdFund",
    "Opex:R&M","Opex:Supplies","Opex:Insurance",
    "Opex:Depreciation"
]
DEPARTMENTS = ["Sales","Ops","HQ"]
DEPT_PARENT = {"Sales":"HQ", "Ops":"HQ", "HQ":None}

def seed_actuals():
    rows = []
    for m in MONTHS:
        sales = 14.0 * 350 * 30
        food = sales * 0.88; bev = sales * 0.12
        cogs_food = food*0.29; cogs_paper = sales*0.02
        labor = sales*0.26; taxes = labor*0.076; bene = labor*0.08
        rent=6500; util=1800; royalty=sales*0.05; ad=sales*0.04
        rnm=sales*0.01; sup=sales*0.005; ins=900; depr=1200
        month=f"2024-{m}"
        def add(a,amt,dept,ccy="USD"): rows.append((a,month,round(amt,2),dept,ccy))
        add("Revenue:Food",food,"Sales"); add("Revenue:Beverage",bev,"Sales")
        add("COGS:Food",cogs_food,"Ops"); add("COGS:Paper",cogs_paper,"Ops")
        add("Opex:Labor",labor,"Ops"); add("Opex:PayrollTaxes",taxes,"HQ"); add("Opex:Benefits",bene,"HQ")
        add("Opex:Rent",rent,"HQ"); add("Opex:Utilities",util,"HQ"); add("Opex:Royalty",royalty,"HQ"); add("Opex:AdFund",ad,"HQ")
        add("Opex:R&M",rnm,"Ops"); add("Opex:Supplies",sup,"Ops"); add("Opex:Insurance",ins,"HQ"); add("Opex:Depreciation",depr,"HQ")
    return pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])
actuals_df = seed_actuals()

budgets: Dict[str, pd.DataFrame] = {}
if "2025:Base" not in budgets:
    b = actuals_df.copy()
    b["month"] = b["month"].str.replace("2024", "2025")
    b["amount"] = (b["amount"] * 1.03).round(2)
    budgets["2025:Base"] = b

fx_rows = []
for y in YEARS:
    for m in MONTHS:
        mm = f"{y}-{m}"
        fx_rows += [("USD",FUNC_CCY,mm,1.0),("EUR",FUNC_CCY,mm,1.1),("INR",FUNC_CCY,mm,0.012)]
fx_rates = pd.DataFrame(fx_rows, columns=["base","quote","month","rate"])

# --------- Headcount store ---------
# Very simple roster table; extend in DB as needed.
roster_cols = ["emp_id","name","dept","start_month","annual_salary","fte","raise_month","raise_pct","benefits_pct","taxes_pct","currency"]
roster_df = pd.DataFrame(columns=roster_cols)

# --------- Helpers ---------
def ensure_period(df: pd.DataFrame, col="month"):
    if col in df.columns and not pd.api.types.is_period_dtype(df[col]):
        df[col] = pd.PeriodIndex(df[col], freq='M')
    return df

def convert_fx(df: pd.DataFrame, func_ccy: str = FUNC_CCY) -> pd.DataFrame:
    if "currency" not in df.columns:
        df = df.assign(currency=func_ccy)
    m = df.merge(fx_rates[fx_rates["quote"]==func_ccy], left_on=["currency","month"], right_on=["base","month"], how="left")
    m["rate"] = m.groupby("currency")["rate"].ffill().bfill().fillna(1.0)
    m["amount_func"] = (m["amount"] * m["rate"]).astype(float)
    return m

def kpi_summary(df: pd.DataFrame) -> Dict:
    d = ensure_period(convert_fx(df.copy()))
    rev = d[d.account_std.str.startswith("Revenue")].groupby('month')['amount_func'].sum()
    cogs = d[d.account_std.str.startswith("COGS")].groupby('month')['amount_func'].sum()
    opex = d[(d.account_std.str.startswith("Opex")) & (d.account_std!="Opex:Depreciation")].groupby('month')['amount_func'].sum()
    depr = d[d.account_std=="Opex:Depreciation"].groupby('month')['amount_func'].sum()
    months = sorted(set(rev.index)|set(cogs.index)|set(opex.index)|set(depr.index))
    def S(s): return pd.Series({m: float(s.get(m,0.0)) for m in months})
    rev,cogs,opex,depr = map(S,[rev,cogs,opex,depr])
    ebitda = rev - cogs - opex
    gm = (rev - cogs) / rev.replace(0,np.nan)
    return {'months':[str(m) for m in months],'revenue':rev,'cogs':cogs,'ebitda':ebitda,'gm':gm.fillna(0)}

def compute_bva(actuals: pd.DataFrame, budget: pd.DataFrame) -> pd.DataFrame:
    a = actuals[["account_std","month","amount","dept","currency"]].copy()
    b = budget[["account_std","month","amount","dept","currency"]].copy().rename(columns={'amount':'amount_budget'})
    m = a.merge(b, on=['account_std','month','dept','currency'], how='outer').fillna(0)
    m_a = convert_fx(m.rename(columns={'amount':'amount_actual'}))
    m_b = convert_fx(m.rename(columns={'amount_budget':'amount'})).rename(columns={'amount_func':'budget_func'})
    m = m_a.join(m_b['budget_func'])
    m['variance'] = m['amount_func'] - m['budget_func']
    def fav(row):
        if row['account_std'].startswith('Revenue'): return 'F' if row['variance']>=0 else 'U'
        return 'F' if row['variance']<=0 else 'U'
    m['FU'] = m.apply(fav, axis=1)
    return m

def ai_variance_narrative(bva_df: pd.DataFrame) -> str:
    d = ensure_period(bva_df.copy())
    totals = d.groupby('account_std')['variance'].sum().sort_values(key=lambda s: s.abs(), ascending=False)
    rev = totals[[i for i in totals.index if i.startswith('Revenue')]].sum() if any(totals.index.str.startswith('Revenue')) else 0.0
    cogs = totals[[i for i in totals.index if i.startswith('COGS')]].sum() if any(totals.index.str.startswith('COGS')) else 0.0
    opex = totals[[i for i in totals.index if i.startswith('Opex')]].sum() if any(totals.index.str.startswith('Opex')) else 0.0
    net = rev - cogs - opex
    top = totals.head(3)
    lines = [f"Overall performance is {'above' if net>=0 else 'below'} plan by {net:,.0f}.",
             f"Revenue Δ: {rev:,.0f}; COGS Δ: {cogs:,.0f}; Opex Δ: {opex:,.0f}.",
             "Top drivers: " + ", ".join([f"{k} ({v:,.0f})" for k,v in top.items()])]
    if cogs>0: lines.append("Costs went up because food prices rose ~15% vs plan and paper costs increased.")
    return " ".join(lines)

def forecast_arima(series: pd.Series, steps:int=12):
    try:
        model = ARIMA(series, order=(1,1,1)).fit()
        f = model.forecast(steps=steps)
        return f
    except Exception:
        growth = (series.diff().tail(3).mean() / max(series.tail(3).mean(),1e-6)) if len(series)>3 else 0.01
        out=[]; last=series.iloc[-1]
        for _ in range(steps):
            last *= (1+float(growth)); out.append(last)
        return pd.Series(out, index=pd.period_range(series.index[-1]+1, periods=steps, freq='M'))

def forecast_ml(series: pd.Series, steps:int=12):
    lag=3; X=[]; y=[]
    for i in range(lag, len(series)):
        X.append(series.values[i-lag:i]); y.append(series.values[i])
    if len(X)<5: return forecast_arima(series, steps)
    reg = LinearRegression().fit(np.array(X), np.array(y))
    hist = list(series.values); preds=[]
    for _ in range(steps):
        x = np.array(hist[-lag:]).reshape(1,-1)
        p = float(reg.predict(x)); preds.append(p); hist.append(p)
    return pd.Series(preds, index=pd.period_range(series.index[-1]+1, periods=steps, freq='M'))

# --------- Headcount math ---------
def headcount_to_budget(roster: pd.DataFrame, fiscal_year:int) -> pd.DataFrame:
    """Expand roster rows into monthly Opex:Labor/Taxes/Benefits for the year."""
    if roster.empty:
        return pd.DataFrame(columns=["account_std","month","amount","dept","currency"])
    out = []
    for _, r in roster.iterrows():
        start = pd.Period(r["start_month"], freq="M")
        for m in [pd.Period(f"{fiscal_year}-{mm}", 'M') for mm in MONTHS]:
            if m < start: continue
            # salary monthly (with raises applied once raise_month is hit)
            annual = float(r["annual_salary"])
            if r.get("raise_month"):
                if str(m) >= str(pd.Period(r["raise_month"], freq='M')):
                    annual *= (1+float(r.get("raise_pct",0.0)))
            monthly = annual/12.0 * float(r.get("fte",1.0))
            taxes = monthly * float(r.get("taxes_pct",0.076))
            benefits = monthly * float(r.get("benefits_pct",0.08))
            month = str(m); dept = r["dept"]; ccy = r.get("currency","USD")
            out.append(("Opex:Labor", month, round(monthly,2), dept, ccy))
            out.append(("Opex:PayrollTaxes", month, round(taxes,2), dept, ccy))
            out.append(("Opex:Benefits", month, round(benefits,2), dept, ccy))
    return pd.DataFrame(out, columns=["account_std","month","amount","dept","currency"])

# --------- Goal Seek solver ---------
def ebitda_of_budget(budget_df: pd.DataFrame) -> float:
    s = kpi_summary(budget_df)
    return float(s["ebitda"].sum())

def apply_levers(df: pd.DataFrame, revenue_pct: float, cogs_pct: float, labor_pct: float) -> pd.DataFrame:
    d = df.copy()
    if revenue_pct != 0:
        mask = d["account_std"].str.startswith("Revenue")
        d.loc[mask,"amount"] = d.loc[mask,"amount"] * (1+revenue_pct)
    if cogs_pct != 0:
        mask = d["account_std"].str.startswith("COGS")
        d.loc[mask,"amount"] = d.loc[mask,"amount"] * (1+cogs_pct)
    if labor_pct != 0:
        mask = d["account_std"]=="Opex:Labor"
        d.loc[mask,"amount"] = d.loc[mask,"amount"] * (1+labor_pct)
    return d

def goal_seek(budget_df: pd.DataFrame, target_ebitda: float, search_range: float = 0.2) -> Tuple[float,float,float,float]:
    """Return (best_ebitda, revenue_pct, cogs_pct, labor_pct) that gets closest to target.
    Coarse grid search for clarity to students.
    """
    base = budget_df.copy()
    best = (float("-inf"), 0.0, 0.0, 0.0)
    best_err = float("inf")
    steps = [ -search_range, -0.1, -0.05, -0.02, 0, 0.02, 0.05, 0.1, search_range ]
    for rp in steps:
        for cp in steps:
            for lp in steps:
                cand = apply_levers(base, rp, cp, lp)
                e = ebitda_of_budget(cand)
                err = abs(e - target_ebitda)
                if err < best_err:
                    best_err = err
                    best = (e, rp, cp, lp)
    return best + (best_err,)

# ------------------- API -------------------
api = FastAPI(title="FP&A Backend v3 (Headcount + Solver)")
api.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

@api.get("/health")
def health(): return {"ok": True}

class LoginReq(BaseModel):
    username: str
    role: Literal["Admin","Analyst","Viewer"]
@api.post("/auth/login")
def login(payload: LoginReq):
    for tok, info in TOKENS.items():
        if info["role"] == payload.role:
            return {"token": tok, "user": payload.username, "role": payload.role}
    raise HTTPException(400, "Unknown role")

# ---- Presets & budgets
class PresetPayload(BaseModel):
    fiscal_year: int = 2026
    overrides: Optional[Dict] = None
    gdp_growth: float = 0.02

def qsr_preset(fy:int, overrides:Optional[Dict], gdp_growth:float=0.0, currency="USD"):
    p = {"avg_ticket":14.0, "daily_txn":380, "open_days_per_month":30,
         "seasonality":[1.00,0.98,1.00,1.02,1.03,1.05,1.08,1.07,1.02,1.01,1.10,1.12],
         "cogs_pct":0.31, "labor_pct":0.26, "rent_fixed":6500.0, "utilities_fixed":1800.0,
         "royalty_pct":0.05, "ad_fund_pct":0.04}
    if overrides: p.update(overrides)
    rows=[]; mg = (1+gdp_growth)**(1/12)-1 if gdp_growth else 0.0
    for i, m in enumerate(MONTHS):
        factor = p["seasonality"][i]
        sales = p["avg_ticket"]*p["daily_txn"]*p["open_days_per_month"]*factor
        sales *= (1+mg)**(i+1)
        food = sales*0.88; bev=sales*0.12
        cogs_food = food*(p["cogs_pct"]*0.9); cogs_paper = sales*(p["cogs_pct"]*0.1)
        labor = sales*p["labor_pct"]; taxes=labor*0.076; bene=labor*0.08
        rent=p["rent_fixed"]; util=p["utilities_fixed"]; royalty=sales*p["royalty_pct"]; ad=sales*p["ad_fund_pct"]
        rnm=sales*0.01; sup=sales*0.005; ins=900; depr=1200; month=f"{fy}-{m}"
        def add(a,amt,dept): rows.append((a,month,round(amt,2),dept,currency))
        add("Revenue:Food",food,"Sales"); add("Revenue:Beverage",bev,"Sales")
        add("COGS:Food",cogs_food,"Ops"); add("COGS:Paper",cogs_paper,"Ops")
        add("Opex:Labor",labor,"Ops"); add("Opex:PayrollTaxes",taxes,"HQ"); add("Opex:Benefits",bene,"HQ")
        add("Opex:Rent",rent,"HQ"); add("Opex:Utilities",util,"HQ"); add("Opex:Royalty",royalty,"HQ"); add("Opex:AdFund",ad,"HQ")
        add("Opex:R&M",rnm,"Ops"); add("Opex:Supplies",sup,"Ops"); add("Opex:Insurance",ins,"HQ"); add("Opex:Depreciation",depr,"HQ")
    return pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])

@api.post("/preset/qsr")
def preset_qsr(payload: PresetPayload, user=Depends(require_role("Admin","Analyst"))):
    df = qsr_preset(payload.fiscal_year, payload.overrides, payload.gdp_growth)
    key = f"{payload.fiscal_year}:Base"
    budgets[key] = df.copy()
    return {"ok": True, "scenario": key, "rows": len(df)}

# ---- Headcount/Roster endpoints
class Emp(BaseModel):
    emp_id: str
    name: str
    dept: Literal["Sales","Ops","HQ"]
    start_month: str   # "YYYY-MM"
    annual_salary: float
    fte: float = 1.0
    raise_month: Optional[str] = None
    raise_pct: float = 0.0
    benefits_pct: float = 0.08
    taxes_pct: float = 0.076
    currency: str = "USD"

@api.get("/headcount/list")
def hc_list(user=Depends(require_role("Admin","Analyst","Viewer"))):
    return {"rows": roster_df.to_dict('records')}

@api.post("/headcount/upsert")
def hc_upsert(emp: Emp, user=Depends(require_role("Admin","Analyst"))):
    global roster_df
    df = pd.DataFrame([emp.dict()])
    # upsert by emp_id
    if not roster_df.empty and (roster_df["emp_id"]==emp.emp_id).any():
        mask = roster_df["emp_id"]==emp.emp_id
        for c in df.columns:
            roster_df.loc[mask, c] = df[c].iloc[0]
    else:
        roster_df = pd.concat([roster_df, df], ignore_index=True)
    return {"ok": True, "count": int(len(roster_df))}

class HCBakeReq(BaseModel):
    fiscal_year: int
    scenario: str = "Base"
    apply: bool = True

@api.post("/headcount/bake_to_budget")
def hc_bake(payload: HCBakeReq, user=Depends(require_role("Admin","Analyst"))):
    """Generate monthly payroll lines from roster and merge into the scenario budget for the year."""
    scen = f"{payload.fiscal_year}:{payload.scenario}" if ":" not in payload.scenario else payload.scenario
    if scen not in budgets: raise HTTPException(404, "scenario not found")
    hc = headcount_to_budget(roster_df, payload.fiscal_year)
    if hc.empty:
        return {"ok": True, "applied_rows": 0}
    # replace only payroll-related rows for that fiscal year
    b = budgets[scen].copy()
    mask_year = b['month'].str.startswith(str(payload.fiscal_year))
    mask_payroll = mask_year & b['account_std'].isin(["Opex:Labor","Opex:PayrollTaxes","Opex:Benefits"])
    b = b.loc[~mask_payroll]
    b = pd.concat([b, hc], ignore_index=True)
    if payload.apply:
        budgets[scen] = b
    return {"ok": True, "applied_rows": len(hc)}

# ---- AutoGen (from trends) and Scenarios/Sensitivity
class AutoGenPayload(BaseModel):
    fiscal_year: int = 2026
    yoy: float = 0.08
    gdp_growth: float = 0.0

@api.post("/budget/autogen")
def budget_autogen(payload: AutoGenPayload, user=Depends(require_role("Admin","Analyst"))):
    d = actuals_df.copy()
    d = ensure_period(d)
    last6 = d[d['month']>= (d['month'].max()-5)].groupby(['account_std','dept','currency'])['amount'].mean()
    months = [pd.Period(f"{payload.fiscal_year}-{m}", 'M') for m in MONTHS]
    rows = []
    mg = (1+payload.yoy)**(1/12) - 1
    gdp_mg = (1+payload.gdp_growth)**(1/12) - 1 if payload.gdp_growth else 0.0
    for i, p in enumerate(months):
        for (acct, dept, ccy), base in last6.items():
            val = float(base) * (1+mg+gdp_mg)**(i+1)
            rows.append((acct, str(p), round(val,2), dept, ccy))
    df = pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])
    key = f"{payload.fiscal_year}:Base"
    budgets[key] = df.copy()
    return {"ok": True, "scenario": key, "rows": len(df)}

class ScenarioClone(BaseModel):
    base: str
    to: str
    pct: float = 0.02
@api.post("/scenario/clone")
def scenario_clone(payload: ScenarioClone, user=Depends(require_role("Admin","Analyst"))):
    if payload.base not in budgets: raise HTTPException(404, "base not found")
    df = budgets[payload.base].copy()
    df["amount"] = (df["amount"]*(1+payload.pct)).round(2)
    budgets[payload.to] = df
    return {"ok": True}

class SensitivityReq(BaseModel):
    scenario: str
    account_pattern: str = "COGS:*"
    start_month: Optional[str] = None
    end_month: Optional[str] = None
    pct: float = 0.1
@api.post("/scenario/sensitivity")
def scenario_sensitivity(payload: SensitivityReq, user=Depends(require_role("Admin","Analyst"))):
    name = payload.scenario
    if name not in budgets: raise HTTPException(404, "scenario not found")
    df = budgets[name].copy()
    mask = df['account_std'].str.startswith(payload.account_pattern[:-2]) if payload.account_pattern.endswith(':*') else (df['account_std']==payload.account_pattern)
    if payload.start_month: mask &= df['month']>=payload.start_month
    if payload.end_month: mask &= df['month']<=payload.end_month
    df.loc[mask,'amount'] = (df.loc[mask,'amount']*(1+payload.pct)).round(2)
    budgets[name] = df
    return {"ok": True, "updated_rows": int(mask.sum())}

# ---- BvA & Forecast
class BvAReq(BaseModel):
    year: int
    scenario: str
@api.post("/bva/analyze")
def bva_analyze(payload: BvAReq, user=Depends(require_role("Admin","Analyst","Viewer"))):
    year = str(payload.year)
    if payload.scenario not in budgets: raise HTTPException(404, "scenario not found")
    a = actuals_df[actuals_df['month'].str.startswith(year)]
    b = budgets[payload.scenario]
    bva = compute_bva(a,b)
    narrative = ai_variance_narrative(bva)
    hot = (bva.assign(abs_var=lambda x: x['variance'].abs()).sort_values('abs_var', ascending=False).head(10)
           [['account_std','dept','month','amount_func','budget_func','variance','FU']])
    return {"summary": narrative, "hotspots": hot.to_dict('records')}

class ForecastReq(BaseModel):
    year: int = 2024
    account: str = "Revenue:Food"
    model: str = "arima"
    months: int = 12
@api.post("/forecast/ai")
def ai_forecast(payload: ForecastReq, user=Depends(require_role("Admin","Analyst","Viewer"))):
    y = str(payload.year)
    a = convert_fx(actuals_df[(actuals_df['month'].str.startswith(y)) & (actuals_df['account_std']==payload.account)])
    if a.empty: raise HTTPException(404, "no actuals for account/year")
    s = ensure_period(a)[['month','amount_func']].set_index('month')['amount_func'].sort_index()
    if payload.model.lower()=="ml":
        f = forecast_ml(s, payload.months)
    else:
        f = forecast_arima(s, payload.months)
    out = pd.DataFrame({'month':[str(i) for i in f.index], 'forecast': [float(x) for x in f.values]})
    return out.to_dict('records')

# ---- Board pack
class BoardReq(BaseModel):
    year: int
    scenario: str
def make_board_pack(summary: str, kpis: Dict[str,str], filepath='board_pack.pptx') -> str:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = 'Executive Summary'
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tb.text_frame.text = summary
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = 'Key KPIs'
    y = 1.5
    for k,v in kpis.items():
        box = s2.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
        box.text_frame.text = f"{k}: {v}"
        y += 0.5
    prs.save(filepath); return filepath
@api.post("/report/boardpack")
def boardpack(payload: BoardReq, user=Depends(require_role("Admin","Analyst"))):
    year = str(payload.year)
    if payload.scenario not in budgets: raise HTTPException(404, "scenario not found")
    a = actuals_df[actuals_df['month'].str.startswith(year)]
    s = kpi_summary(a)
    summary = f"Auto board summary for {year}. GM avg {(s['gm'].mean()*100):.1f}%."
    path = make_board_pack(summary, {
        "Revenue (YTD)": f"${s['revenue'].sum():,.0f}",
        "EBITDA (YTD)": f"${s['ebitda'].sum():,.0f}",
        "GM% (avg)": f"{(s['gm'].mean()*100):.1f}%"
    })
    return {"ok": True, "file": path}

# ---- QuickBooks mapper stub
class QBImportJSON(BaseModel):
    rows: List[Dict]
@api.post("/qb/import_json")
def qb_import_json(payload: QBImportJSON, user=Depends(require_role("Admin","Analyst"))):
    df = pd.DataFrame(payload.rows)
    mapping = {
        "Sales - Food": "Revenue:Food",
        "Sales - Beverage": "Revenue:Beverage",
        "COGS - Food": "COGS:Food",
        "COGS - Paper": "COGS:Paper",
        "Wages": "Opex:Labor",
        "Payroll Taxes": "Opex:PayrollTaxes",
        "Benefits": "Opex:Benefits",
        "Rent": "Opex:Rent",
        "Utilities": "Opex:Utilities",
        "Royalty": "Opex:Royalty",
        "Advertising": "Opex:AdFund",
        "Repairs": "Opex:R&M",
        "Supplies": "Opex:Supplies",
        "Insurance": "Opex:Insurance",
        "Depreciation": "Opex:Depreciation",
    }
    for col in ["Date","Account","Amount"]:
        if col not in df.columns: raise HTTPException(400, f"Missing {col} column")
    if 'Dept' not in df.columns: df['Dept'] = 'HQ'
    if 'Currency' not in df.columns: df['Currency'] = FUNC_CCY
    df["account_std"] = df["Account"].map(mapping).fillna(df["Account"])
    df["month"] = pd.to_datetime(df["Date"]).dt.to_period("M").astype(str)
    df = df.rename(columns={"Dept":"dept","Amount":"amount","Currency":"currency"})[["account_std","month","amount","dept","currency"]]
    global actuals_df
    actuals_df = pd.concat([actuals_df, df], ignore_index=True)
    return {"ok": True, "imported": len(df)}

@api.post("/qb/import_csv")
def qb_import_csv(file: UploadFile = File(...), user=Depends(require_role("Admin","Analyst"))):
    content = file.file.read().decode("utf-8")
    reader = csv.DictReader(io.StringIO(content))
    rows = list(reader)
    return qb_import_json(QBImportJSON(rows=rows), user)

# ---- Excel endpoints
@api.post("/excel/retrieve")
def excel_retrieve(payload: dict = Body(...), user=Depends(require_role("Admin","Analyst","Viewer"))):
    year = str(payload.get('year', 2025)); scen_short = payload.get('scenario', 'Base')
    scenario = f"{year}:{scen_short}" if ":" not in scen_short else scen_short
    df = budgets.get(scenario)
    if df is None: raise HTTPException(404, 'Scenario not found')
    out = df[df['month'].str.startswith(year)].copy()
    accts = payload.get('accounts')
    if accts:
        mask = pd.Series(False, index=out.index)
        for pat in accts:
            if pat.endswith(':*'): mask |= out['account_std'].str.startswith(pat[:-2])
            else: mask |= (out['account_std']==pat)
        out = out[mask]
    return out.sort_values(['account_std','dept','month']).to_dict('records')

@api.post("/excel/submit")
def excel_submit(payload: dict = Body(...), user=Depends(require_role("Admin","Analyst"))):
    scen = payload.get('scenario','2025:Base'); rows = pd.DataFrame(payload.get('rows', []))
    if scen not in budgets:
        budgets[scen] = rows
        return {'ok': True, 'created': scen, 'rows': len(rows)}
    df = budgets[scen]; key = ['account_std','dept','month']
    budgets[scen] = (pd.concat([df, rows]).drop_duplicates(subset=key, keep='last').reset_index(drop=True))
    return {'ok': True, 'updated': len(rows)}

# ---- Solver endpoint
class SolveReq(BaseModel):
    scenario: str
    target_ebitda: float
    search_range: float = 0.2
    apply: bool = False
@api.post("/solver/goal_seek")
def solver_goal_seek(payload: SolveReq, user=Depends(require_role("Admin","Analyst"))):
    scen = payload.scenario
    if scen not in budgets: raise HTTPException(404, "scenario not found")
    base = budgets[scen]
    best_ebitda, rp, cp, lp, err = goal_seek(base, payload.target_ebitda, payload.search_range)
    result = {"best_ebitda": round(best_ebitda,2), "revenue_pct": rp, "cogs_pct": cp, "labor_pct": lp, "abs_error": round(err,2)}
    if payload.apply:
        budgets[scen] = apply_levers(base, rp, cp, lp)
        result["applied_to"] = scen
    return result
