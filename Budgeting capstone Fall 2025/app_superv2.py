"""
app_super.py — The single-file FP&A backend + dashboard (all-in-one)

This file is deliberately monolithic so students can run ONE Python file and get:
- FastAPI API (budgets, scenarios, forecasting, BvA, versions, Excel endpoints)
- Plotly Dash executive dashboard at /dashboard
- Auth & roles (Admin/Analyst/Viewer) via Bearer token
- GDP-driven presets, multi-currency, departments
- AI forecasting (ARIMA/ML), AI narrative explanations
- QuickBooks import (stub), Board pack PPT export
- Version control & audit log (in-memory)

Run:
  pip install fastapi uvicorn pandas numpy statsmodels scikit-learn python-pptx plotly dash python-multipart
  uvicorn app_super:api --reload --port 8000
Then open:
  API docs → http://127.0.0.1:8000/docs
  Dashboard → http://127.0.0.1:8000/dashboard

Note: In-memory data for teaching. Swap with a DB in production.
"""
from __future__ import annotations
import os, io, csv, math
from typing import Dict, Optional, List, Literal
import numpy as np
import pandas as pd

from fastapi import FastAPI, Body, HTTPException, Depends, UploadFile, File, Header
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from starlette.middleware.wsgi import WSGIMiddleware

# Forecasting libs
from statsmodels.tsa.arima.model import ARIMA
from sklearn.linear_model import LinearRegression

# Board pack PPT
from pptx import Presentation
from pptx.util import Inches

# Dash/Plotly
from dash import Dash, dcc, html, Input, Output
import plotly.express as px
import plotly.graph_objects as go

# -------------------
# Auth & Roles
# -------------------
security = HTTPBearer(auto_error=False)
TOKENS = {
    "admin-token": {"user":"alice","role":"Admin"},
    "analyst-token": {"user":"bob","role":"Analyst"},
    "viewer-token": {"user":"cathy","role":"Viewer"},
}
def get_current_user(creds: Optional[HTTPAuthorizationCredentials] = Depends(security)):
    if creds is None or not creds.credentials:
        raise HTTPException(401, "Missing Bearer token")
    info = TOKENS.get(creds.credentials)
    if not info:
        raise HTTPException(403, "Invalid token")
    return info
def require_role(*roles: str):
    def checker(user=Depends(get_current_user)):
        if user["role"] not in roles:
            raise HTTPException(403, f"Requires role: {', '.join(roles)}")
        return user
    return checker

# -------------------
# In-memory DATA
# -------------------
FUNC_CCY = os.getenv("FUNC_CURRENCY", "USD")
YEARS = [2024, 2025, 2026]
MONTHS = [f"{m:02d}" for m in range(1,13)]
ACCOUNTS = [
    "Revenue:Food","Revenue:Beverage",
    "COGS:Food","COGS:Paper",
    "Opex:Labor","Opex:PayrollTaxes","Opex:Benefits",
    "Opex:Rent","Opex:Utilities","Opex:Royalty","Opex:AdFund",
    "Opex:R&M","Opex:Supplies","Opex:Insurance",
    "Opex:Depreciation"
]
DEPARTMENTS = ["Sales","Ops","HQ"]
DEPT_PARENT = {"Sales":"HQ", "Ops":"HQ", "HQ":None}

def seed_actuals():
    rows = []
    for m in MONTHS:
        sales = 14.0 * 350 * 30
        food = sales * 0.88; bev = sales * 0.12
        cogs_food = food*0.29; cogs_paper = sales*0.02
        labor = sales*0.26; taxes = labor*0.076; bene = labor*0.08
        rent=6500; util=1800; royalty=sales*0.05; ad=sales*0.04
        rnm=sales*0.01; sup=sales*0.005; ins=900; depr=1200
        month=f"2024-{m}"
        def add(a,amt,dept,ccy="USD"): rows.append((a,month,round(amt,2),dept,ccy))
        add("Revenue:Food",food,"Sales"); add("Revenue:Beverage",bev,"Sales")
        add("COGS:Food",cogs_food,"Ops"); add("COGS:Paper",cogs_paper,"Ops")
        add("Opex:Labor",labor,"Ops"); add("Opex:PayrollTaxes",taxes,"HQ"); add("Opex:Benefits",bene,"HQ")
        add("Opex:Rent",rent,"HQ"); add("Opex:Utilities",util,"HQ"); add("Opex:Royalty",royalty,"HQ"); add("Opex:AdFund",ad,"HQ")
        add("Opex:R&M",rnm,"Ops"); add("Opex:Supplies",sup,"Ops"); add("Opex:Insurance",ins,"HQ"); add("Opex:Depreciation",depr,"HQ")
    return pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])
actuals_df = seed_actuals()

# Budgets keyed by "YYYY:Scenario"
budgets: Dict[str, pd.DataFrame] = {}
versions: Dict[str, Dict[str, pd.DataFrame]] = {}   # scenario -> name -> df
audit_log: List[Dict] = []

# Seed a default 2025 Base from actuals uplift
if "2025:Base" not in budgets:
    b = actuals_df.copy()
    b["month"] = b["month"].str.replace("2024", "2025")
    b["amount"] = (b["amount"] * 1.03).round(2)
    budgets["2025:Base"] = b

# FX (monthly). For demo, USD=1, EUR=1.1, INR=0.012
fx_rows = []
for y in YEARS:
    for m in MONTHS:
        mm = f"{y}-{m}"
        fx_rows += [("USD",FUNC_CCY,mm,1.0),("EUR",FUNC_CCY,mm,1.1),("INR",FUNC_CCY,mm,0.012)]
fx_rates = pd.DataFrame(fx_rows, columns=["base","quote","month","rate"])

# ------------------- Helpers -------------------
def ensure_period(df: pd.DataFrame, col="month"):
    if col in df.columns and not pd.api.types.is_period_dtype(df[col]):
        df[col] = pd.PeriodIndex(df[col], freq='M')
    return df

def convert_fx(df: pd.DataFrame, func_ccy: str = FUNC_CCY) -> pd.DataFrame:
    if "currency" not in df.columns:
        df = df.assign(currency=func_ccy)
    m = df.merge(fx_rates[fx_rates["quote"]==func_ccy], left_on=["currency","month"], right_on=["base","month"], how="left")
    m["rate"] = m.groupby("currency")["rate"].ffill().bfill().fillna(1.0)
    m["amount_func"] = (m["amount"] * m["rate"]).astype(float)
    return m

def kpi_summary(df: pd.DataFrame) -> Dict:
    d = ensure_period(convert_fx(df.copy()))
    rev = d[d.account_std.str.startswith("Revenue")].groupby('month')['amount_func'].sum()
    cogs = d[d.account_std.str.startswith("COGS")].groupby('month')['amount_func'].sum()
    opex = d[(d.account_std.str.startswith("Opex")) & (d.account_std!="Opex:Depreciation")].groupby('month')['amount_func'].sum()
    depr = d[d.account_std=="Opex:Depreciation"].groupby('month')['amount_func'].sum()
    months = sorted(set(rev.index)|set(cogs.index)|set(opex.index)|set(depr.index))
    def S(s): return pd.Series({m: float(s.get(m,0.0)) for m in months})
    rev,cogs,opex,depr = map(S,[rev,cogs,opex,depr])
    ebitda = rev - cogs - opex
    gm = (rev - cogs) / rev.replace(0,np.nan)
    return {'months':[str(m) for m in months],'revenue':rev,'cogs':cogs,'ebitda':ebitda,'gm':gm.fillna(0)}

def compute_bva(actuals: pd.DataFrame, budget: pd.DataFrame) -> pd.DataFrame:
    a = actuals[["account_std","month","amount","dept","currency"]].copy()
    b = budget[["account_std","month","amount","dept","currency"]].copy().rename(columns={'amount':'amount_budget'})
    m = a.merge(b, on=['account_std','month','dept','currency'], how='outer').fillna(0)
    m_a = convert_fx(m.rename(columns={'amount':'amount_actual'}))
    m_b = convert_fx(m.rename(columns={'amount_budget':'amount'})).rename(columns={'amount_func':'budget_func'})
    m = m_a.join(m_b['budget_func'])
    m['variance'] = m['amount_func'] - m['budget_func']
    def fav(row):
        if row['account_std'].startswith('Revenue'): return 'F' if row['variance']>=0 else 'U'
        return 'F' if row['variance']<=0 else 'U'
    m['FU'] = m.apply(fav, axis=1)
    return m

def ai_variance_narrative(bva_df: pd.DataFrame) -> str:
    d = ensure_period(bva_df.copy())
    totals = d.groupby('account_std')['variance'].sum().sort_values(key=lambda s: s.abs(), ascending=False)
    rev = totals[[i for i in totals.index if i.startswith('Revenue')]].sum() if any(totals.index.str.startswith('Revenue')) else 0.0
    cogs = totals[[i for i in totals.index if i.startswith('COGS')]].sum() if any(totals.index.str.startswith('COGS')) else 0.0
    opex = totals[[i for i in totals.index if i.startswith('Opex')]].sum() if any(totals.index.str.startswith('Opex')) else 0.0
    net = rev - cogs - opex
    top = totals.head(3)
    lines = [f"Overall performance is {'above' if net>=0 else 'below'} plan by {net:,.0f}.",
             f"Revenue Δ: {rev:,.0f}; COGS Δ: {cogs:,.0f}; Opex Δ: {opex:,.0f}.",
             "Top drivers: " + ", ".join([f"{k} ({v:,.0f})" for k,v in top.items()])]
    if cogs>0: lines.append("Costs went up because food prices rose ~15% vs plan and paper costs increased.")
    return " ".join(lines)

def forecast_arima(series: pd.Series, steps:int=12):
    try:
        model = ARIMA(series, order=(1,1,1)).fit()
        f = model.forecast(steps=steps)
        return f
    except Exception:
        growth = (series.diff().tail(3).mean() / max(series.tail(3).mean(),1e-6)) if len(series)>3 else 0.01
        out=[]; last=series.iloc[-1]
        for _ in range(steps):
            last *= (1+float(growth)); out.append(last)
        return pd.Series(out, index=pd.period_range(series.index[-1]+1, periods=steps, freq='M'))

def forecast_ml(series: pd.Series, steps:int=12):
    lag=3; X=[]; y=[]
    for i in range(lag, len(series)):
        X.append(series.values[i-lag:i]); y.append(series.values[i])
    if len(X)<5: return forecast_arima(series, steps)
    reg = LinearRegression().fit(np.array(X), np.array(y))
    hist = list(series.values); preds=[]
    for _ in range(steps):
        x = np.array(hist[-lag:]).reshape(1,-1)
        p = float(reg.predict(x)); preds.append(p); hist.append(p)
    return pd.Series(preds, index=pd.period_range(series.index[-1]+1, periods=steps, freq='M'))

def record_audit(action:str, detail:Dict, user:Dict):
    audit_log.append({"action":action, "detail":detail, "user":user["user"], "role":user["role"]})

# -------------------
# FastAPI
# -------------------
api = FastAPI(title="FP&A SUPER (single-file)")
api.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

@api.get("/health")
def health(): return {"ok": True}

# ---- Auth endpoints
class LoginReq(BaseModel):
    username: str
    role: Literal["Admin","Analyst","Viewer"]
@api.post("/auth/login")
def login(payload: LoginReq):
    for tok, info in TOKENS.items():
        if info["role"] == payload.role:
            return {"token": tok, "user": payload.username, "role": payload.role}
    raise HTTPException(400, "Unknown role")

# ---- Presets & budgets
class PresetPayload(BaseModel):
    fiscal_year: int = 2026
    overrides: Optional[Dict] = None
    gdp_growth: float = 0.02

def qsr_preset(fy:int, overrides:Optional[Dict], gdp_growth:float=0.0, currency="USD"):
    p = {"avg_ticket":14.0, "daily_txn":380, "open_days_per_month":30,
         "seasonality":[1.00,0.98,1.00,1.02,1.03,1.05,1.08,1.07,1.02,1.01,1.10,1.12],
         "cogs_pct":0.31, "labor_pct":0.26, "rent_fixed":6500.0, "utilities_fixed":1800.0,
         "royalty_pct":0.05, "ad_fund_pct":0.04}
    if overrides: p.update(overrides)
    rows=[]; mg = (1+gdp_growth)**(1/12)-1 if gdp_growth else 0.0
    for i, m in enumerate(MONTHS):
        factor = p["seasonality"][i]
        sales = p["avg_ticket"]*p["daily_txn"]*p["open_days_per_month"]*factor
        sales *= (1+mg)**(i+1)
        food = sales*0.88; bev=sales*0.12
        cogs_food = food*(p["cogs_pct"]*0.9); cogs_paper = sales*(p["cogs_pct"]*0.1)
        labor = sales*p["labor_pct"]; taxes=labor*0.076; bene=labor*0.08
        rent=p["rent_fixed"]; util=p["utilities_fixed"]; royalty=sales*p["royalty_pct"]; ad=sales*p["ad_fund_pct"]
        rnm=sales*0.01; sup=sales*0.005; ins=900; depr=1200; month=f"{fy}-{m}"
        def add(a,amt,dept): rows.append((a,month,round(amt,2),dept,currency))
        add("Revenue:Food",food,"Sales"); add("Revenue:Beverage",bev,"Sales")
        add("COGS:Food",cogs_food,"Ops"); add("COGS:Paper",cogs_paper,"Ops")
        add("Opex:Labor",labor,"Ops"); add("Opex:PayrollTaxes",taxes,"HQ"); add("Opex:Benefits",bene,"HQ")
        add("Opex:Rent",rent,"HQ"); add("Opex:Utilities",util,"HQ"); add("Opex:Royalty",royalty,"HQ"); add("Opex:AdFund",ad,"HQ")
        add("Opex:R&M",rnm,"Ops"); add("Opex:Supplies",sup,"Ops"); add("Opex:Insurance",ins,"HQ"); add("Opex:Depreciation",depr,"HQ")
    return pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])

@api.post("/preset/qsr")
def preset_qsr(payload: PresetPayload, user=Depends(require_role("Admin","Analyst"))):
    df = qsr_preset(payload.fiscal_year, payload.overrides, payload.gdp_growth)
    key = f"{payload.fiscal_year}:Base"
    budgets[key] = df.copy()
    record_audit("preset_qsr", {"scenario":key}, user)
    return {"ok": True, "scenario": key, "rows": len(df)}

class AutoGenPayload(BaseModel):
    fiscal_year: int = 2026
    yoy: float = 0.08
    gdp_growth: float = 0.0
@api.post("/budget/autogen")
def budget_autogen(payload: AutoGenPayload, user=Depends(require_role("Admin","Analyst"))):
    d = actuals_df.copy()
    d = ensure_period(d)
    last6 = d[d['month']>= (d['month'].max()-5)].groupby(['account_std','dept','currency'])['amount'].mean()
    months = [pd.Period(f"{payload.fiscal_year}-{m}", 'M') for m in MONTHS]
    rows = []
    mg = (1+payload.yoy)**(1/12) - 1
    gdp_mg = (1+payload.gdp_growth)**(1/12) - 1 if payload.gdp_growth else 0.0
    for i, p in enumerate(months):
        for (acct, dept, ccy), base in last6.items():
            val = float(base) * (1+mg+gdp_mg)**(i+1)
            rows.append((acct, str(p), round(val,2), dept, ccy))
    df = pd.DataFrame(rows, columns=["account_std","month","amount","dept","currency"])
    key = f"{payload.fiscal_year}:Base"
    budgets[key] = df.copy()
    record_audit("budget_autogen", {"scenario":key}, user)
    return {"ok": True, "scenario": key, "rows": len(df)}

# ---- Scenarios, sensitivity, versions
class ScenarioClone(BaseModel):
    base: str
    to: str
    pct: float = 0.02
@api.post("/scenario/clone")
def scenario_clone(payload: ScenarioClone, user=Depends(require_role("Admin","Analyst"))):
    if payload.base not in budgets: raise HTTPException(404, "base not found")
    df = budgets[payload.base].copy()
    df["amount"] = (df["amount"]*(1+payload.pct)).round(2)
    budgets[payload.to] = df
    record_audit("scenario_clone", {"from":payload.base, "to":payload.to, "pct":payload.pct}, user)
    return {"ok": True}

class SensitivityReq(BaseModel):
    scenario: str
    account_pattern: str = "COGS:*"
    start_month: Optional[str] = None
    end_month: Optional[str] = None
    pct: float = 0.1
@api.post("/scenario/sensitivity")
def scenario_sensitivity(payload: SensitivityReq, user=Depends(require_role("Admin","Analyst"))):
    name = payload.scenario
    if name not in budgets: raise HTTPException(404, "scenario not found")
    df = budgets[name].copy()
    mask = df['account_std'].str.startswith(payload.account_pattern[:-2]) if payload.account_pattern.endswith(':*') else (df['account_std']==payload.account_pattern)
    if payload.start_month: mask &= df['month']>=payload.start_month
    if payload.end_month: mask &= df['month']<=payload.end_month
    df.loc[mask,'amount'] = (df.loc[mask,'amount']*(1+payload.pct)).round(2)
    budgets[name] = df
    record_audit("scenario_sensitivity", {"scenario":name, "pattern":payload.account_pattern, "pct":payload.pct}, user)
    return {"ok": True, "updated_rows": int(mask.sum())}

class VersionReq(BaseModel):
    scenario: str
    name: str
@api.post("/versions/save")
def versions_save(payload: VersionReq, user=Depends(require_role("Admin","Analyst"))):
    if payload.scenario not in budgets: raise HTTPException(404, "scenario not found")
    versions.setdefault(payload.scenario, {})[payload.name] = budgets[payload.scenario].copy()
    record_audit("versions_save", {"scenario":payload.scenario, "name":payload.name}, user)
    return {"ok": True, "versions": list(versions[payload.scenario].keys())}

@api.get("/versions/list")
def versions_list(scenario: str, user=Depends(require_role("Admin","Analyst","Viewer"))):
    return {"scenario": scenario, "versions": list(versions.get(scenario, {}).keys())}

class RestoreReq(BaseModel):
    scenario: str
    name: str
@api.post("/versions/restore")
def versions_restore(payload: RestoreReq, user=Depends(require_role("Admin","Analyst"))):
    snap = versions.get(payload.scenario, {}).get(payload.name)
    if snap is None: raise HTTPException(404, "version not found")
    budgets[payload.scenario] = snap.copy()
    record_audit("versions_restore", {"scenario":payload.scenario, "name":payload.name}, user)
    return {"ok": True}

@api.get("/audit/logs")
def audit_logs(user=Depends(require_role("Admin","Analyst","Viewer"))):
    return {"events": audit_log[-200:]}

# ---- BvA & Forecast
class BvAReq(BaseModel):
    year: int
    scenario: str
@api.post("/bva/analyze")
def bva_analyze(payload: BvAReq, user=Depends(require_role("Admin","Analyst","Viewer"))):
    year = str(payload.year)
    if payload.scenario not in budgets: raise HTTPException(404, "scenario not found")
    a = actuals_df[actuals_df['month'].str.startswith(year)]
    b = budgets[payload.scenario]
    bva = compute_bva(a,b)
    narrative = ai_variance_narrative(bva)
    hot = (bva.assign(abs_var=lambda x: x['variance'].abs()).sort_values('abs_var', ascending=False).head(10)
           [['account_std','dept','month','amount_func','budget_func','variance','FU']])
    return {"summary": narrative, "hotspots": hot.to_dict('records')}

class ForecastReq(BaseModel):
    year: int = 2024
    account: str = "Revenue:Food"
    model: str = "arima"
    months: int = 12
@api.post("/forecast/ai")
def ai_forecast(payload: ForecastReq, user=Depends(require_role("Admin","Analyst","Viewer"))):
    y = str(payload.year)
    a = convert_fx(actuals_df[(actuals_df['month'].str.startswith(y)) & (actuals_df['account_std']==payload.account)])
    if a.empty: raise HTTPException(404, "no actuals for account/year")
    s = ensure_period(a)[['month','amount_func']].set_index('month')['amount_func'].sort_index()
    if payload.model.lower()=="ml":
        f = forecast_ml(s, payload.months)
    else:
        f = forecast_arima(s, payload.months)
    out = pd.DataFrame({'month':[str(i) for i in f.index], 'forecast': [float(x) for x in f.values]})
    return out.to_dict('records')

# ---- Board pack PPT
class BoardReq(BaseModel):
    year: int
    scenario: str
def make_board_pack(summary: str, kpis: Dict[str,str], filepath='board_pack.pptx') -> str:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = 'Executive Summary'
    tb = s1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tb.text_frame.text = summary
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = 'Key KPIs'
    y = 1.5
    for k,v in kpis.items():
        box = s2.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.4))
        box.text_frame.text = f"{k}: {v}"
        y += 0.5
    prs.save(filepath); return filepath
@api.post("/report/boardpack")
def boardpack(payload: BoardReq, user=Depends(require_role("Admin","Analyst"))):
    year = str(payload.year)
    if payload.scenario not in budgets: raise HTTPException(404, "scenario not found")
    a = actuals_df[actuals_df['month'].str.startswith(year)]
    s = kpi_summary(a)
    summary = f"Auto board summary for {year}. GM avg {(s['gm'].mean()*100):.1f}%."
    path = make_board_pack(summary, {
        "Revenue (YTD)": f"${s['revenue'].sum():,.0f}",
        "EBITDA (YTD)": f"${s['ebitda'].sum():,.0f}",
        "GM% (avg)": f"{(s['gm'].mean()*100):.1f}%"
    })
    record_audit("boardpack", {"file":path}, user)
    return {"ok": True, "file": path}

# ---- QuickBooks mapper stub
class QBImportJSON(BaseModel):
    rows: List[Dict]
@api.post("/qb/import_json")
def qb_import_json(payload: QBImportJSON, user=Depends(require_role("Admin","Analyst"))):
    df = pd.DataFrame(payload.rows)
    mapping = {
        "Sales - Food": "Revenue:Food",
        "Sales - Beverage": "Revenue:Beverage",
        "COGS - Food": "COGS:Food",
        "COGS - Paper": "COGS:Paper",
        "Wages": "Opex:Labor",
        "Payroll Taxes": "Opex:PayrollTaxes",
        "Benefits": "Opex:Benefits",
        "Rent": "Opex:Rent",
        "Utilities": "Opex:Utilities",
        "Royalty": "Opex:Royalty",
        "Advertising": "Opex:AdFund",
        "Repairs": "Opex:R&M",
        "Supplies": "Opex:Supplies",
        "Insurance": "Opex:Insurance",
        "Depreciation": "Opex:Depreciation",
    }
    if 'Date' not in df.columns: raise HTTPException(400, "Missing Date column")
    if 'Account' not in df.columns: raise HTTPException(400, "Missing Account column")
    if 'Dept' not in df.columns: df['Dept'] = 'HQ'
    if 'Currency' not in df.columns: df['Currency'] = FUNC_CCY
    if 'Amount' not in df.columns: raise HTTPException(400, "Missing Amount column")
    df["account_std"] = df["Account"].map(mapping).fillna(df["Account"])
    df["month"] = pd.to_datetime(df["Date"]).dt.to_period("M").astype(str)
    df = df.rename(columns={"Dept":"dept","Amount":"amount","Currency":"currency"})[["account_std","month","amount","dept","currency"]]
    global actuals_df
    actuals_df = pd.concat([actuals_df, df], ignore_index=True)
    record_audit("qb_import", {"rows":len(df)}, user)
    return {"ok": True, "imported": len(df)}

@api.post("/qb/import_csv")
def qb_import_csv(file: UploadFile = File(...), user=Depends(require_role("Admin","Analyst"))):
    content = file.file.read().decode("utf-8")
    reader = csv.DictReader(io.StringIO(content))
    rows = list(reader)
    return qb_import_json(QBImportJSON(rows=rows), user)

# ---- Excel-style endpoints (retrieve/submit/drill)
@api.post("/excel/retrieve")
def excel_retrieve(payload: dict = Body(...), user=Depends(require_role("Admin","Analyst","Viewer"))):
    year = str(payload.get('year', 2025)); scen_short = payload.get('scenario', 'Base')
    scenario = f"{year}:{scen_short}" if ":" not in scen_short else scen_short
    df = budgets.get(scenario)
    if df is None: raise HTTPException(404, 'Scenario not found')
    out = df[df['month'].str.startswith(year)].copy()
    accts = payload.get('accounts')
    if accts:
        mask = pd.Series(False, index=out.index)
        for pat in accts:
            if pat.endswith(':*'): mask |= out['account_std'].str.startswith(pat[:-2])
            else: mask |= (out['account_std']==pat)
        out = out[mask]
    return out.sort_values(['account_std','dept','month']).to_dict('records')

@api.post("/excel/submit")
def excel_submit(payload: dict = Body(...), user=Depends(require_role("Admin","Analyst"))):
    scen = payload.get('scenario','2025:Base'); rows = pd.DataFrame(payload.get('rows', []))
    if scen not in budgets:
        budgets[scen] = rows
        record_audit("excel_submit_create", {"scenario":scen, "rows":len(rows)}, user)
        return {'ok': True, 'created': scen, 'rows': len(rows)}
    df = budgets[scen]; key = ['account_std','dept','month']
    budgets[scen] = (pd.concat([df, rows]).drop_duplicates(subset=key, keep='last').reset_index(drop=True))
    record_audit("excel_submit_update", {"scenario":scen, "rows":len(rows)}, user)
    return {'ok': True, 'updated': len(rows)}

@api.get("/excel/expand_children")
def excel_expand_children(member: str, axis: str = 'Department', user=Depends(require_role("Admin","Analyst","Viewer"))):
    if axis=='Department':
        if member=='HQ': return {'member':'HQ','children':['Sales','Ops']}
        return {'member':member,'children':[]}
    return {'member':member,'children':[]}

# -------------------
# Embedded Dash Dashboard
# -------------------
dash_app = Dash(__name__)
dash_app.title = "Executive Dashboard"

def chart_trend(summary: Dict) -> go.Figure:
    df = pd.DataFrame({'month': summary['months'],
                       'Revenue': summary['revenue'].values,
                       'COGS': summary['cogs'].values,
                       'EBITDA': summary['ebitda'].values})
    fig = px.line(df, x='month', y=['Revenue','COGS','EBITDA'], markers=True)
    fig.update_layout(title='Revenue / COGS / EBITDA Trend', legend_title=None)
    return fig

def chart_yoy(a_prev: pd.DataFrame, a_cur: pd.DataFrame, label1: str, label2: str) -> go.Figure:
    def total_by_account(d):
        dd = convert_fx(d.copy()); return dd.groupby('account_std', as_index=False)['amount_func'].sum()
    m = total_by_account(a_prev).merge(total_by_account(a_cur), on='account_std', how='outer', suffixes=(f'_{label1}', f'_{label2}')).fillna(0)
    fig = go.Figure()
    fig.add_bar(x=m['account_std'], y=m[f'amount_{label1}'], name=label1)
    fig.add_bar(x=m['account_std'], y=m[f'amount_{label2}'], name=label2)
    fig.update_layout(barmode='group', title=f'YoY Comparison {label1} vs {label2}')
    return fig

def chart_waterfall(bva_df: pd.DataFrame) -> go.Figure:
    m = (bva_df.groupby('account_std', as_index=False)['variance'].sum().sort_values('variance'))
    fig = go.Figure(go.Waterfall(name='Variance', orientation='v', measure=['relative']*len(m), x=m['account_std'], y=m['variance']))
    fig.update_layout(title='Variance Waterfall (Actual - Budget)')
    return fig

def chart_heatmap(bva_df: pd.DataFrame) -> go.Figure:
    d = ensure_period(bva_df.copy())
    piv = d.pivot_table(index='account_std', columns='month', values='variance', aggfunc='sum').fillna(0)
    fig = px.imshow(piv, aspect='auto', labels=dict(x='Month', y='Account', color='Variance'))
    fig.update_layout(title='Variance Heatmap')
    return fig

dash_app.layout = html.Div([
    html.H2("Executive Dashboard (Single-file)"),
    html.Div([
        html.Label('Year'), dcc.Dropdown(options=[{'label':y,'value':y} for y in YEARS], value=2025, id='dd-year', style={'width':'120px'}),
        html.Label('Scenario'), dcc.Dropdown(options=[{'label':s,'value':s.split(':')[-1]} for s in ["2025:Base","2025:Best","2025:Worst"]], value='Base', id='dd-scenario', style={'width':'160px'}),
        html.Button("Refresh", id='btn-refresh', style={'marginLeft':'12px'})
    ], style={'display':'flex','gap':'1rem','alignItems':'center'}),
    html.Div(id='kpi-cards', style={'display':'flex','gap':'2rem','marginTop':'1rem'}),
    dcc.Graph(id='trend'),
    dcc.Graph(id='yoy'),
    html.Div([dcc.Graph(id='waterfall'), dcc.Graph(id='heatmap')], style={'display':'grid','gridTemplateColumns':'1fr 1fr','gap':'1rem'}),
])

@dash_app.callback(
    Output('kpi-cards','children'),
    Output('trend','figure'),
    Output('yoy','figure'),
    Output('waterfall','figure'),
    Output('heatmap','figure'),
    Input('dd-year','value'),
    Input('dd-scenario','value'),
    Input('btn-refresh','n_clicks')
)
def update_dashboard(year, scenario_short, _n):
    scenario = f"{year}:{scenario_short}"
    if scenario not in budgets:
        # create defaults Best/Worst from Base if missing
        if f"{year}:Base" in budgets:
            budgets[f"{year}:Best"] = budgets[f"{year}:Base"].assign(amount=lambda x: (x['amount']*1.03).round(2))
            budgets[f"{year}:Worst"] = budgets[f"{year}:Base"].assign(amount=lambda x: (x['amount']*0.97).round(2))
    a = actuals_df[actuals_df['month'].str.startswith(str(year))]
    s = kpi_summary(a)
    cards = [
        html.Div([html.Div('Revenue', className='label'), html.H3(f"${s['revenue'].sum():,.0f}")], className='card'),
        html.Div([html.Div('EBITDA', className='label'), html.H3(f"${s['ebitda'].sum():,.0f}")], className='card'),
        html.Div([html.Div('GM %', className='label'), html.H3(f"{(s['gm'].mean()*100):.1f}%")], className='card')
    ]
    trend_fig = chart_trend(s)
    prev = year-1; a_prev = actuals_df[actuals_df['month'].str.startswith(str(prev))]; a_cur = actuals_df[actuals_df['month'].str.startswith(str(year))]
    yoy_fig = chart_yoy(a_prev, a_cur, str(prev), str(year))
    b = budgets.get(scenario, budgets.get(f"{year}:Base", list(budgets.values())[0]))
    bva = compute_bva(a_cur, b)
    wf = chart_waterfall(bva); hm = chart_heatmap(bva)
    return cards, trend_fig, yoy_fig, wf, hm

api.mount("/dashboard", WSGIMiddleware(dash_app.server))
